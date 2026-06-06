import os, json, re, math, logging, io
from datetime import datetime, timedelta, timezone
from functools import wraps
from collections import defaultdict
from pathlib import Path

from flask import Flask, request, jsonify, send_from_directory, g, redirect
from flask_cors import CORS
import jwt, bcrypt

# ── PostgreSQL & Serverless Storage ───────────────────────────────────────────
import psycopg2
from psycopg2 import pool
from psycopg2.extras import RealDictCursor
import cloudinary
import cloudinary.uploader
import cloudinary.api

def _utcnow(): return datetime.now(timezone.utc)

# ── Optional NLP & Extraction Libs ────────────────────────────────────────────
try:
    import fitz;                    HAS_PDF  = True
except ImportError:                 HAS_PDF  = False
try:
    from docx import Document as _DocxDoc; HAS_DOCX = True
except ImportError:                 HAS_DOCX = False
try:
    from pptx import Presentation;  HAS_PPTX = True
except ImportError:                 HAS_PPTX = False
try:
    from sentence_transformers import SentenceTransformer
    import numpy as np
    _AI = SentenceTransformer('all-MiniLM-L6-v2')
    HAS_ST = True; HAS_NP = True
    print("[AI] SentenceTransformer loaded ✓")
except Exception as _e:
    HAS_ST = False
    print(f"[AI] Fallback TF-IDF ({_e})")
    try:
        import numpy as np; HAS_NP = True
    except ImportError:     HAS_NP = False

# ── Configuration ─────────────────────────────────────────────────────────────
BASE_DIR   = os.path.dirname(os.path.abspath(__file__))
DATABASE_URL = os.environ.get("DATABASE_URL")
SECRET     = os.environ.get("SECRET_KEY", "gsu-faculty-repo-secret-key-2024!!")
JWT_DAYS   = int(os.environ.get("JWT_DAYS", 7))
MAX_BYTES  = int(os.environ.get("MAX_MB", 50)) * 1024 * 1024
PORT       = int(os.environ.get("PORT", 5000))
ALLOWED    = {".pdf",".docx",".pptx",".txt",".xlsx",".png",".jpg",".jpeg"}

# Cloudinary Setup
cloudinary.config(
    cloud_name = os.environ.get('CLOUDINARY_CLOUD_NAME'),
    api_key = os.environ.get('CLOUDINARY_API_KEY'),
    api_secret = os.environ.get('CLOUDINARY_API_SECRET'),
    secure = True
)

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger("GSURepo")

app = Flask(__name__, static_folder=str(BASE_DIR), static_url_path="")
app.config["MAX_CONTENT_LENGTH"] = MAX_BYTES
CORS(app, resources={r"/api/*": {"origins": "*"}}, supports_credentials=True, allow_headers=["Content-Type","Authorization"])

# ── Database Connection Pooling (Neon) ────────────────────────────────────────
try:
    db_pool = psycopg2.pool.ThreadedConnectionPool(1, 15, dsn=DATABASE_URL)
    log.info("Neon PostgreSQL connection pool initialized.")
except Exception as e:
    log.error(f"Failed to connect to Neon PostgreSQL: {e}")
    db_pool = None

def get_db():
    if "db_conn" not in g:
        g.db_conn = db_pool.getconn()
    return g.db_conn

def execute_query(query, params=None, fetch=None, commit=False):
    conn = get_db()
    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute(query, params)
            if commit:
                conn.commit()
            if fetch == "one":
                return cur.fetchone()
            elif fetch == "all":
                return cur.fetchall()
            elif fetch == "id":
                return cur.fetchone()['id']
    except Exception as e:
        conn.rollback()
        log.error(f"DB Error: {e}")
        raise e

@app.teardown_appcontext
def close_db(e=None):
    conn = g.pop("db_conn", None)
    if conn:
        db_pool.putconn(conn)

def init_db():
    with app.app_context():
        execute_query("""
        CREATE TABLE IF NOT EXISTS users (
            id SERIAL PRIMARY KEY,
            name TEXT NOT NULL, email TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL, role TEXT DEFAULT 'faculty',
            department TEXT DEFAULT '', bio TEXT DEFAULT '',
            active INTEGER DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            last_login TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS documents (
            id SERIAL PRIMARY KEY,
            title TEXT NOT NULL, description TEXT DEFAULT '',
            filename TEXT NOT NULL, original_name TEXT NOT NULL,
            file_type TEXT NOT NULL, file_size INTEGER DEFAULT 0,
            course_code TEXT DEFAULT '', academic_level TEXT DEFAULT '',
            resource_type TEXT DEFAULT '', academic_year TEXT DEFAULT '',
            uploader_id INTEGER NOT NULL REFERENCES users(id), 
            is_public INTEGER DEFAULT 1,
            text_content TEXT DEFAULT '', embedding TEXT DEFAULT '',
            download_count INTEGER DEFAULT 0, view_count INTEGER DEFAULT 0,
            cloudinary_url TEXT DEFAULT '', cloudinary_public_id TEXT DEFAULT '',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS tags (
            id SERIAL PRIMARY KEY,
            document_id INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE, 
            tag TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS bookmarks (
            id SERIAL PRIMARY KEY,
            user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE, 
            document_id INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(user_id, document_id)
        );
        CREATE TABLE IF NOT EXISTS comments (
            id SERIAL PRIMARY KEY,
            document_id INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE, 
            user_id INTEGER NOT NULL REFERENCES users(id),
            content TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS search_logs (
            id SERIAL PRIMARY KEY,
            user_id INTEGER, query TEXT NOT NULL, result_count INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS activity_log (
            id SERIAL PRIMARY KEY,
            user_id INTEGER, action TEXT NOT NULL, document_id INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE INDEX IF NOT EXISTS idx_d_uploader ON documents(uploader_id);
        CREATE INDEX IF NOT EXISTS idx_d_course   ON documents(course_code);
        CREATE INDEX IF NOT EXISTS idx_t_doc      ON tags(document_id);
        CREATE INDEX IF NOT EXISTS idx_t_tag      ON tags(tag);
        CREATE INDEX IF NOT EXISTS idx_bm         ON bookmarks(user_id);
        CREATE INDEX IF NOT EXISTS idx_cmt        ON comments(document_id);
        """, commit=True)
        
        admin_exists = execute_query("SELECT id FROM users WHERE role='admin' LIMIT 1", fetch="one")
        if not admin_exists:
            pw = bcrypt.hashpw(b"Admin@1234", bcrypt.gensalt()).decode()
            execute_query(
                "INSERT INTO users (name,email,password,role,department) VALUES (%s,%s,%s,%s,%s)",
                ("System Admin","admin@gsu.edu.ng", pw,"admin","Computer Science"), commit=True
            )
            log.info("Default admin created: admin@gsu.edu.ng / Admin@1234")

# ── Auth Helpers ──────────────────────────────────────────────────────────────
def make_token(uid, role):
    t = jwt.encode({"sub":uid,"role":role, "exp": _utcnow()+timedelta(days=JWT_DAYS), "iat": _utcnow()}, SECRET, algorithm="HS256")
    return t if isinstance(t, str) else t.decode()

def require_auth(f):
    @wraps(f)
    def w(*a, **kw):
        auth = request.headers.get("Authorization","")
        if not auth.startswith("Bearer "): return jsonify({"error":"Missing token"}), 401
        try:
            d = jwt.decode(auth.split(" ",1)[1], SECRET, algorithms=["HS256"])
            g.user_id = d["sub"]; g.user_role = d["role"]
        except jwt.ExpiredSignatureError: return jsonify({"error":"Token expired"}), 401
        except Exception: return jsonify({"error":"Invalid token"}), 401
        
        # Verify user is active
        u = execute_query("SELECT active FROM users WHERE id=%s", (g.user_id,), fetch="one")
        if not u or u['active'] == 0:
            return jsonify({"error":"Account deactivated"}), 403
            
        return f(*a, **kw)
    return w

def require_admin(f):
    @wraps(f)
    @require_auth
    def w(*a, **kw):
        if g.user_role != "admin": return jsonify({"error":"Admin only"}), 403
        return f(*a, **kw)
    return w

# ── In-Memory Serverless Text Extraction ──────────────────────────────────────
def extract_text_from_bytes(file_bytes, ext):
    try:
        if ext == ".pdf"  and HAS_PDF:
            return " ".join(p.get_text() for p in fitz.open(stream=file_bytes, filetype="pdf"))[:50000]
        if ext == ".docx" and HAS_DOCX:
            return " ".join(p.text for p in _DocxDoc(io.BytesIO(file_bytes)).paragraphs)[:50000]
        if ext == ".pptx" and HAS_PPTX:
            prs = Presentation(io.BytesIO(file_bytes))
            return " ".join(s.text for sl in prs.slides for s in sl.shapes if hasattr(s,"text"))[:50000]
        if ext == ".txt":
            return file_bytes.decode('utf-8', errors='ignore')[:50000]
    except Exception as e:
        log.warning(f"extract_text error: {e}")
    return ""

# ── AI / NLP Classification & Semantic Engine ─────────────────────────────────
COURSE_KW = {
    "CSC101":["programming","algorithm","introduction","computer","basic"],
    "CSC201":["data structure","stack","queue","linked list","tree","binary"],
    "CSC301":["database","sql","relational","normalization","query"],
    "CSC302":["operating system","process","memory","scheduling","thread"],
    "CSC401":["artificial intelligence","machine learning","neural","deep learning"],
    "CSC402":["network","protocol","tcp","ip","routing","socket"],
    "CSC403":["software engineering","sdlc","agile","design pattern","uml"],
    "CSC404":["computer graphics","rendering","opengl","pixel","3d"],
    "CSC501":["research","methodology","thesis","dissertation"],
    "MAT101":["mathematics","calculus","algebra","differential","integral"],
    "STA101":["statistics","probability","regression","hypothesis"],
}
RES_KW = {
    "Lecture Note":  ["lecture","note","introduction","overview","chapter"],
    "Research Paper":["abstract","methodology","conclusion","references","journal"],
    "Assignment":    ["assignment","exercise","problem","submit","question"],
    "Past Question": ["past question","examination","exam","test","quiz"],
    "Textbook":      ["textbook","edition","publisher","isbn"],
    "Presentation":  ["slide","presentation","powerpoint"],
    "Tutorial":      ["tutorial","guide","step by step","how to"],
    "Project":       ["project","final year","implementation","system design"],
}
LVL_KW = {
    "100 Level":  ["100 level","year one","freshman","first year"],
    "200 Level":  ["200 level","year two","sophomore"],
    "300 Level":  ["300 level","year three","junior"],
    "400 Level":  ["400 level","year four","senior","final year"],
    "Postgraduate":["postgraduate","masters","phd","msc"],
}
STOPS = {"the","a","an","and","or","but","in","on","at","to","for","of","with",
         "is","are","was","were","be","been","have","has","had","do","does",
         "did","will","would","could","should","that","this","those","it","its",
         "we","our","they","their","which","from","by","as","not","no","all"}

def _score(t, km):
    s = {l: sum(t.count(k) for k in ks) for l,ks in km.items()}
    s = {l:v for l,v in s.items() if v>0}
    return max(s,key=s.get) if s else None

def auto_classify(text, fname):
    t = (text+" "+fname).lower()
    return {"course_code":  _score(t,COURSE_KW) or "",
            "resource_type":_score(t,RES_KW) or "General",
            "academic_level":_score(t,LVL_KW) or ""}

def extract_tags(text, n=12):
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
    freq  = defaultdict(int)
    for w in words:
        if w not in STOPS: freq[w] += 1
    for i in range(len(words)-1):
        a,b = words[i],words[i+1]
        if a not in STOPS and b not in STOPS: freq[f"{a} {b}"] += 2
    return [w for w,_ in sorted(freq.items(),key=lambda x:x[1],reverse=True)[:n]]

def get_embedding(text):
    if HAS_ST and text.strip():
        try: return _AI.encode(text[:5000], convert_to_numpy=True).tolist()
        except: pass
    return None

def cos_np(a,b):
    if not HAS_NP: return 0.0
    va,vb = np.array(a),np.array(b)
    return float(np.dot(va,vb)/(np.linalg.norm(va)*np.linalg.norm(vb)+1e-9))

def tfidf(corpus):
    N   = len(corpus)
    tok = [re.findall(r'\b\w+\b', d.lower()) for d in corpus]
    df  = defaultdict(int)
    for ts in tok:
        for t in set(ts): df[t] += 1
    vecs = []
    for ts in tok:
        tf = defaultdict(int)
        for t in ts: tf[t] += 1
        vecs.append({t:(c/(len(ts) or 1))*(math.log((N+1)/(df[t]+1))+1)
                     for t,c in tf.items()})
    return vecs

def cos_dict(v1,v2):
    keys = set(v1)&set(v2)
    if not keys: return 0.0
    dot = sum(v1[k]*v2[k] for k in keys)
    return dot/(math.sqrt(sum(x*x for x in v1.values()))*
                math.sqrt(sum(x*x for x in v2.values()))+1e-9)

# ── Search Engine ─────────────────────────────────────────────────────────────
def search_docs(query, filters, uid=None):
    execute_query("INSERT INTO search_logs (user_id,query) VALUES (%s,%s)", (uid, query), commit=True)
    sql = ("SELECT d.*, u.name as uploader_name, STRING_AGG(t.tag, ',') as tags "
           "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
           "LEFT JOIN tags t ON t.document_id=d.id WHERE d.is_public=1")
    p = []
    if filters.get("course"):  sql += " AND d.course_code=%s";   p.append(filters["course"])
    if filters.get("year"):    sql += " AND d.academic_year=%s"; p.append(filters["year"])
    if filters.get("type"):    sql += " AND d.resource_type=%s"; p.append(filters["type"])
    if filters.get("level"):   sql += " AND d.academic_level=%s";p.append(filters["level"])
    if filters.get("author") and str(filters["author"]).isdigit():
        sql += " AND d.uploader_id=%s"; p.append(int(filters["author"]))
    
    sql += " GROUP BY d.id, u.name"
    rows = execute_query(sql, p, fetch="all")
    docs = [dict(r) for r in rows]
    
    if not query.strip():
        docs.sort(key=lambda x:x["created_at"],reverse=True); return docs[:50]
    
    ql = query.lower()
    if HAS_ST:
        try:
            qe = _AI.encode(query,convert_to_numpy=True).tolist()
            scored = []
            for d in docs:
                kw=0
                for f in ["title","description","course_code","tags"]:
                    val=str(d.get(f,"") or "").lower()
                    if ql in val: kw+=3
                    for w in ql.split():
                        if w in val: kw+=1
                sem=0.0
                if d.get("embedding"):
                    try: sem=cos_np(qe,json.loads(d["embedding"]))*10
                    except: pass
                tot=sem+kw
                if tot>0: d["relevance"]=round(tot,4); scored.append(d)
            return sorted(scored,key=lambda x:x["relevance"],reverse=True)[:50]
        except Exception as e: log.warning(f"Semantic Search Error: {e}")
        
    # TF-IDF Fallback
    corpus=[f"{d['title']} {d['description']} {d['text_content']}" for d in docs]
    if not corpus: return []
    vecs=tfidf(corpus+[query]); qv=vecs[-1]; scored=[]
    for i,d in enumerate(docs):
        sim=cos_dict(vecs[i],qv)
        if ql in str(d.get("title","")).lower(): sim+=0.5
        if sim>0.01: d["relevance"]=round(sim,4); scored.append(d)
    return sorted(scored,key=lambda x:x["relevance"],reverse=True)[:50]

def get_recs(doc_id, n=6):
    target = execute_query("SELECT * FROM documents WHERE id=%s", (doc_id,), fetch="one")
    if not target: return []
    target = dict(target)
    
    cands = execute_query(
        "SELECT d.*, STRING_AGG(t.tag, ',') as tags FROM documents d "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.id!=%s AND d.is_public=1 GROUP BY d.id",
        (doc_id,), fetch="all"
    )
    cands = [dict(r) for r in cands]
    if not cands: return []
    
    scores=[]
    if HAS_ST and target.get("embedding"):
        try:
            te=json.loads(target["embedding"])
            for c in cands:
                s=(2 if c["course_code"]==target["course_code"] and c["course_code"] else 0)
                s+=(1 if c["academic_level"]==target["academic_level"] and c["academic_level"] else 0)
                if c.get("embedding"):
                    try: s+=cos_np(te,json.loads(c["embedding"]))*5
                    except: pass
                scores.append((c,s))
        except: scores=[]
    if not scores:
        corpus=[f"{d['title']} {d['description']} {d['text_content']}" for d in cands]
        tt=f"{target['title']} {target['description']} {target['text_content']}"
        vecs=tfidf(corpus+[tt]); tv=vecs[-1]
        for i,c in enumerate(cands):
            sim=cos_dict(vecs[i],tv)
            if c["course_code"]==target["course_code"] and c["course_code"]: sim+=0.3
            scores.append((c,sim))
            
    scores.sort(key=lambda x:x[1],reverse=True)
    return [dict(d,relevance=round(s,4)) for d,s in scores[:n]]

# ═══════════════════════════════════════════════════════════════════════════════
#  API ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/",defaults={"path":""})
@app.route("/<path:path>")
def spa(path):
    if path.startswith("api/") or path.startswith("uploads/"): return jsonify({"error":"Not found"}),404
    return send_from_directory(str(BASE_DIR),"index.html")

@app.route("/api/health")
def health():
    doc_count = execute_query("SELECT COUNT(*) as count FROM documents", fetch="one")["count"]
    usr_count = execute_query("SELECT COUNT(*) as count FROM users", fetch="one")["count"]
    return jsonify({
        "status":"ok","ai":"sentence-transformers" if HAS_ST else "tfidf",
        "documents": doc_count, "users": usr_count, "timestamp": _utcnow().isoformat()
    })

# ── Authentication ────────────────────────────────────────────────────────────
@app.route("/api/auth/register",methods=["POST"])
def register():
    d=request.json or {}
    name=(d.get("name") or "").strip(); email=(d.get("email") or "").strip().lower()
    pw=(d.get("password") or "").strip(); dept=(d.get("department") or "Computer Science").strip()
    if not all([name,email,pw]): return jsonify({"error":"All fields required"}),400
    if len(pw)<6: return jsonify({"error":"Password min 6 chars"}),400
    
    if execute_query("SELECT id FROM users WHERE email=%s", (email,), fetch="one"):
        return jsonify({"error":"Email already registered"}),409
        
    h = bcrypt.hashpw(pw.encode(),bcrypt.gensalt()).decode()
    uid = execute_query(
        "INSERT INTO users (name,email,password,role,department) VALUES (%s,%s,%s,%s,%s) RETURNING id",
        (name,email,h,"faculty",dept), fetch="id", commit=True
    )
    return jsonify({"token":make_token(uid,"faculty"), "user":{"id":uid,"name":name,"email":email,"role":"faculty","department":dept,"bio":""}}),201

@app.route("/api/auth/login",methods=["POST"])
def login():
    d=request.json or {}
    email=(d.get("email") or "").strip().lower(); pw=(d.get("password") or "").strip()
    if not email or not pw: return jsonify({"error":"Email and password required"}),400
    
    u = execute_query("SELECT * FROM users WHERE email=%s AND active=1", (email,), fetch="one")
    if not u or not bcrypt.checkpw(pw.encode(),u["password"].encode()):
        return jsonify({"error":"Invalid credentials"}),401
        
    execute_query("UPDATE users SET last_login=CURRENT_TIMESTAMP WHERE id=%s", (u["id"],), commit=True)
    execute_query("INSERT INTO activity_log (user_id,action) VALUES (%s,%s)", (u["id"],"login"), commit=True)
    
    return jsonify({"token":make_token(u["id"],u["role"]), "user":{k:u[k] for k in ("id","name","email","role","department","bio")}})

@app.route("/api/auth/me")
@require_auth
def get_me():
    u = execute_query("SELECT id,name,email,role,department,bio,created_at,last_login FROM users WHERE id=%s", (g.user_id,), fetch="one")
    return jsonify(dict(u)) if u else (jsonify({"error":"Not found"}),404)

@app.route("/api/auth/change-password",methods=["PUT"])
@require_auth
def change_pw():
    d=request.json or {}
    old=(d.get("old_password") or "").strip(); new=(d.get("new_password") or "").strip()
    if not old or not new or len(new)<6: return jsonify({"error":"Invalid input"}),400
    
    u = execute_query("SELECT password FROM users WHERE id=%s", (g.user_id,), fetch="one")
    if not bcrypt.checkpw(old.encode(),u["password"].encode()): return jsonify({"error":"Old password wrong"}),401
    
    execute_query("UPDATE users SET password=%s WHERE id=%s", (bcrypt.hashpw(new.encode(),bcrypt.gensalt()).decode(),g.user_id), commit=True)
    return jsonify({"message":"Password changed"})

@app.route("/api/profile",methods=["PUT"])
@require_auth
def update_profile():
    d=request.json or {}; sets,params=[],[]
    for f in ["name","department","bio"]:
        if f in d: sets.append(f"{f}=%s"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    params.append(g.user_id)
    execute_query(f"UPDATE users SET {','.join(sets)} WHERE id=%s", params, commit=True)
    u = execute_query("SELECT id,name,email,role,department,bio FROM users WHERE id=%s", (g.user_id,), fetch="one")
    return jsonify({"message":"Updated","user":dict(u)})

# ── Documents & Files ─────────────────────────────────────────────────────────
def _strip(docs):
    for d in docs: d.pop("text_content",None); d.pop("embedding",None)
    return docs

@app.route("/api/documents")
@require_auth
def list_docs():
    query=request.args.get("q","").strip()
    page=max(1,int(request.args.get("page",1))); per=20
    filters={k:request.args.get(k,"") for k in ("course","year","type","level","author")}
    
    if query or any(filters.values()):
        results = search_docs(query,filters,g.user_id); total=len(results)
        items = _strip(results[(page-1)*per:page*per])
    else:
        total = execute_query("SELECT COUNT(*) as count FROM documents WHERE is_public=1", fetch="one")["count"]
        rows = execute_query(
            "SELECT d.*, u.name as uploader_name, STRING_AGG(t.tag, ',') as tags "
            "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
            "LEFT JOIN tags t ON t.document_id=d.id WHERE d.is_public=1 "
            "GROUP BY d.id, u.name ORDER BY d.created_at DESC LIMIT %s OFFSET %s",
            (per, (page-1)*per), fetch="all"
        )
        items = _strip([dict(r) for r in rows])
    return jsonify({"documents":items,"total":total,"page":page,"per_page":per})

@app.route("/api/documents/<int:doc_id>")
@require_auth
def get_doc(doc_id):
    doc = execute_query(
        "SELECT d.*, u.name as uploader_name, u.email as uploader_email, STRING_AGG(t.tag, ',') as tags "
        "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.id=%s GROUP BY d.id, u.name, u.email",
        (doc_id,), fetch="one"
    )
    if not doc: return jsonify({"error":"Not found"}),404
    
    execute_query("UPDATE documents SET view_count=view_count+1 WHERE id=%s", (doc_id,), commit=True)
    execute_query("INSERT INTO activity_log (user_id,action,document_id) VALUES (%s,%s,%s)", (g.user_id,"view",doc_id), commit=True)
    
    bm = execute_query("SELECT id FROM bookmarks WHERE user_id=%s AND document_id=%s", (g.user_id,doc_id), fetch="one")
    d = dict(doc); d.pop("text_content",None); d.pop("embedding",None)
    d["bookmarked"] = bool(bm)
    return jsonify(d)

@app.route("/api/documents/upload",methods=["POST"])
@require_auth
def upload_doc():
    if "file" not in request.files: return jsonify({"error":"No file"}),400
    file = request.files["file"]
    if not file.filename: return jsonify({"error":"Empty filename"}),400
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED: return jsonify({"error":f"Type {ext} not allowed"}),400
    
    file_bytes = file.read()
    size = len(file_bytes)
    if size > MAX_BYTES: return jsonify({"error":"File too large"}),413
    
    title=request.form.get("title",file.filename)
    desc=request.form.get("description","")
    cc=request.form.get("course_code",""); al=request.form.get("academic_level","")
    rt=request.form.get("resource_type",""); ay=request.form.get("academic_year","")
    pub=int(request.form.get("is_public",1))
    
    text = extract_text_from_bytes(file_bytes, ext)
    if not (cc and al and rt):
        cls = auto_classify(text,file.filename)
        cc = cc or cls["course_code"]; al = al or cls["academic_level"]; rt = rt or cls["resource_type"]
        
    emb = get_embedding(f"{title} {desc} {text[:3000]}")
    tags = extract_tags(f"{title} {desc} {text}")
    
    # Cloudinary Streaming
    file.seek(0)
    c_res_type = "image" if ext in {".png", ".jpg", ".jpeg"} else "raw"
    try:
        upload_result = cloudinary.uploader.upload(
            file, resource_type=c_res_type, folder="gsu_repo", use_filename=True, unique_filename=True
        )
        c_url = upload_result.get("secure_url")
        c_pid = upload_result.get("public_id")
    except Exception as e:
        log.error(f"Cloudinary upload failed: {e}")
        return jsonify({"error": "Failed to upload file to cloud storage."}), 500

    did = execute_query(
        "INSERT INTO documents (title,description,filename,original_name,file_type,file_size,"
        "course_code,academic_level,resource_type,academic_year,uploader_id,text_content,embedding,is_public,cloudinary_url,cloudinary_public_id) "
        "VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) RETURNING id",
        (title,desc,c_pid,file.filename,ext,size,cc,al,rt,ay,g.user_id,text,json.dumps(emb) if emb else "",pub,c_url,c_pid),
        fetch="id", commit=True
    )
    
    for tag in tags: 
        execute_query("INSERT INTO tags (document_id,tag) VALUES (%s,%s)", (did,tag), commit=True)
    execute_query("INSERT INTO activity_log (user_id,action,document_id) VALUES (%s,%s,%s)", (g.user_id,"upload",did), commit=True)
    
    return jsonify({"message":"Uploaded","id":did,"tags":tags,
                    "classification":{"course_code":cc,"academic_level":al,"resource_type":rt}}),201

@app.route("/api/documents/<int:doc_id>",methods=["PUT"])
@require_auth
def update_doc(doc_id):
    doc = execute_query("SELECT * FROM documents WHERE id=%s", (doc_id,), fetch="one")
    if not doc: return jsonify({"error":"Not found"}),404
    if dict(doc)["uploader_id"] != g.user_id and g.user_role != "admin": return jsonify({"error":"Forbidden"}),403
    
    d = request.json or {}; sets,params=[],[]
    for f in ["title","description","course_code","academic_level","resource_type","academic_year","is_public"]:
        if f in d: sets.append(f"{f}=%s"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    
    sets.append("updated_at=CURRENT_TIMESTAMP"); params.append(doc_id)
    execute_query(f"UPDATE documents SET {','.join(sets)} WHERE id=%s", params, commit=True)
    return jsonify({"message":"Updated"})

@app.route("/api/documents/<int:doc_id>",methods=["DELETE"])
@require_auth
def delete_doc(doc_id):
    doc = execute_query("SELECT * FROM documents WHERE id=%s", (doc_id,), fetch="one")
    if not doc: return jsonify({"error":"Not found"}),404
    d = dict(doc)
    if d["uploader_id"] != g.user_id and g.user_role != "admin": return jsonify({"error":"Forbidden"}),403
    
    # Destroy from Cloudinary
    if d.get("cloudinary_public_id"):
        c_res_type = "image" if d["file_type"] in {".png", ".jpg", ".jpeg"} else "raw"
        try: cloudinary.uploader.destroy(d["cloudinary_public_id"], resource_type=c_res_type)
        except Exception as e: log.error(f"Failed to delete from Cloudinary: {e}")

    execute_query("DELETE FROM documents WHERE id=%s", (doc_id,), commit=True)
    return jsonify({"message":"Deleted"})

@app.route("/api/documents/<int:doc_id>/download")
@require_auth
def download_doc(doc_id):
    doc = execute_query("SELECT cloudinary_url, original_name FROM documents WHERE id=%s", (doc_id,), fetch="one")
    if not doc: return jsonify({"error":"Not found"}),404
    
    execute_query("UPDATE documents SET download_count=download_count+1 WHERE id=%s", (doc_id,), commit=True)
    execute_query("INSERT INTO activity_log (user_id,action,document_id) VALUES (%s,%s,%s)", (g.user_id,"download",doc_id), commit=True)
    
    return redirect(doc["cloudinary_url"])

@app.route("/api/documents/<int:doc_id>/recommendations")
@require_auth
def recs(doc_id):
    r = get_recs(doc_id)
    for x in r: x.pop("text_content",None); x.pop("embedding",None)
    return jsonify({"recommendations":r})

@app.route("/api/documents/my")
@require_auth
def my_docs():
    rows = execute_query(
        "SELECT d.*, u.name as uploader_name, STRING_AGG(t.tag, ',') as tags "
        "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.uploader_id=%s "
        "GROUP BY d.id, u.name ORDER BY d.created_at DESC", (g.user_id,), fetch="all"
    )
    return jsonify({"documents":_strip([dict(r) for r in rows])})

# ── Bookmarks ─────────────────────────────────────────────────────────────────
@app.route("/api/bookmarks")
@require_auth
def bookmarks():
    rows = execute_query(
        "SELECT d.*, u.name as uploader_name, STRING_AGG(t.tag, ',') as tags "
        "FROM bookmarks b JOIN documents d ON b.document_id=d.id "
        "LEFT JOIN users u ON d.uploader_id=u.id LEFT JOIN tags t ON t.document_id=d.id "
        "WHERE b.user_id=%s GROUP BY d.id, b.created_at, u.name ORDER BY b.created_at DESC", 
        (g.user_id,), fetch="all"
    )
    return jsonify({"documents":_strip([dict(r) for r in rows])})

@app.route("/api/bookmarks/<int:doc_id>",methods=["POST"])
@require_auth
def toggle_bm(doc_id):
    ex = execute_query("SELECT id FROM bookmarks WHERE user_id=%s AND document_id=%s", (g.user_id,doc_id), fetch="one")
    if ex: 
        execute_query("DELETE FROM bookmarks WHERE user_id=%s AND document_id=%s", (g.user_id,doc_id), commit=True)
        return jsonify({"bookmarked":False})
    execute_query("INSERT INTO bookmarks (user_id,document_id) VALUES (%s,%s)", (g.user_id,doc_id), commit=True)
    return jsonify({"bookmarked":True})

# ── Comments ──────────────────────────────────────────────────────────────────
@app.route("/api/documents/<int:doc_id>/comments")
@require_auth
def get_comments(doc_id):
    rows = execute_query(
        "SELECT c.*, u.name as user_name FROM comments c "
        "JOIN users u ON c.user_id=u.id WHERE c.document_id=%s ORDER BY c.created_at ASC",
        (doc_id,), fetch="all"
    )
    return jsonify({"comments":[dict(r) for r in rows]})

@app.route("/api/documents/<int:doc_id>/comments",methods=["POST"])
@require_auth
def add_comment(doc_id):
    d=request.json or {}; content=(d.get("content") or "").strip()
    if not content: return jsonify({"error":"Empty comment"}),400
    
    cid = execute_query(
        "INSERT INTO comments (document_id,user_id,content) VALUES (%s,%s,%s) RETURNING id",
        (doc_id,g.user_id,content), fetch="id", commit=True
    )
    row = execute_query(
        "SELECT c.*, u.name as user_name FROM comments c JOIN users u ON c.user_id=u.id WHERE c.id=%s",
        (cid,), fetch="one"
    )
    return jsonify(dict(row)),201

@app.route("/api/comments/<int:cid>",methods=["DELETE"])
@require_auth
def del_comment(cid):
    row = execute_query("SELECT * FROM comments WHERE id=%s", (cid,), fetch="one")
    if not row: return jsonify({"error":"Not found"}),404
    if dict(row)["user_id"] != g.user_id and g.user_role != "admin": return jsonify({"error":"Forbidden"}),403
    execute_query("DELETE FROM comments WHERE id=%s", (cid,), commit=True)
    return jsonify({"message":"Deleted"})

# ── Metadata & Dashboards ─────────────────────────────────────────────────────
@app.route("/api/tags/popular")
@require_auth
def pop_tags():
    rows = execute_query("SELECT tag, COUNT(*) as cnt FROM tags GROUP BY tag ORDER BY cnt DESC LIMIT 40", fetch="all")
    return jsonify({"tags":[dict(r) for r in rows]})

@app.route("/api/faculty")
@require_auth
def faculty():
    rows = execute_query(
        "SELECT u.id, u.name, u.department, u.bio, COUNT(d.id) as doc_count "
        "FROM users u LEFT JOIN documents d ON d.uploader_id=u.id "
        "WHERE u.active=1 GROUP BY u.id ORDER BY doc_count DESC", fetch="all"
    )
    return jsonify({"faculty":[dict(u) for u in rows]})

@app.route("/api/admin/stats")
@require_admin
def admin_stats():
    return jsonify({
        "total_documents": execute_query("SELECT COUNT(*) as count FROM documents", fetch="one")["count"],
        "total_users":     execute_query("SELECT COUNT(*) as count FROM users", fetch="one")["count"],
        "total_downloads": execute_query("SELECT COALESCE(SUM(download_count),0) as count FROM documents", fetch="one")["count"],
        "total_views":     execute_query("SELECT COALESCE(SUM(view_count),0) as count FROM documents", fetch="one")["count"],
        "total_searches":  execute_query("SELECT COUNT(*) as count FROM search_logs", fetch="one")["count"],
        "total_bookmarks": execute_query("SELECT COUNT(*) as count FROM bookmarks", fetch="one")["count"],
        "total_comments":  execute_query("SELECT COUNT(*) as count FROM comments", fetch="one")["count"],
        "by_type":    [dict(r) for r in execute_query("SELECT resource_type, COUNT(*) as cnt FROM documents GROUP BY resource_type ORDER BY cnt DESC", fetch="all")],
        "by_course":  [dict(r) for r in execute_query("SELECT course_code, COUNT(*) as cnt FROM documents WHERE course_code!='' GROUP BY course_code ORDER BY cnt DESC LIMIT 10", fetch="all")],
        "by_level":   [dict(r) for r in execute_query("SELECT academic_level, COUNT(*) as cnt FROM documents WHERE academic_level!='' GROUP BY academic_level ORDER BY cnt DESC", fetch="all")],
        "recent_uploads": [dict(r) for r in execute_query("SELECT d.id, d.title, u.name as uploader_name, d.created_at, d.file_type, d.download_count FROM documents d JOIN users u ON d.uploader_id=u.id ORDER BY d.created_at DESC LIMIT 10", fetch="all")],
        "daily_uploads":  [dict(r) for r in execute_query("SELECT DATE(created_at) as date, COUNT(*) as cnt FROM documents GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30", fetch="all")],
        "daily_searches": [dict(r) for r in execute_query("SELECT DATE(created_at) as date, COUNT(*) as cnt FROM search_logs GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30", fetch="all")],
    })

@app.route("/api/admin/users")
@require_admin
def admin_users():
    rows = execute_query(
        "SELECT u.id, u.name, u.email, u.role, u.department, u.active, u.created_at, u.last_login, "
        "COUNT(d.id) as doc_count FROM users u LEFT JOIN documents d ON d.uploader_id=u.id "
        "GROUP BY u.id ORDER BY u.created_at DESC", fetch="all"
    )
    return jsonify({"users":[dict(u) for u in rows]})

@app.route("/api/admin/users/<int:uid>",methods=["PUT"])
@require_admin
def admin_upd_user(uid):
    d=request.json or {}; sets,params=[],[]
    for f in ["role","active","department","name"]:
        if f in d: sets.append(f"{f}=%s"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    params.append(uid)
    execute_query(f"UPDATE users SET {','.join(sets)} WHERE id=%s", params, commit=True)
    return jsonify({"message":"Updated"})

@app.route("/api/admin/users/<int:uid>",methods=["DELETE"])
@require_admin
def admin_del_user(uid):
    if uid==g.user_id: return jsonify({"error":"Cannot deactivate yourself"}),400
    execute_query("UPDATE users SET active=0 WHERE id=%s", (uid,), commit=True)
    return jsonify({"message":"Deactivated"})

@app.route("/api/admin/documents")
@require_admin
def admin_all_docs():
    rows = execute_query(
        "SELECT d.*, u.name as uploader_name, STRING_AGG(t.tag, ',') as tags "
        "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id GROUP BY d.id, u.name ORDER BY d.created_at DESC", 
        fetch="all"
    )
    return jsonify({"documents":_strip([dict(r) for r in rows])})

@app.route("/api/admin/reports/usage")
@require_admin
def usage_report():
    return jsonify({
        "top_documents": [dict(r) for r in execute_query("SELECT id,title,download_count,view_count FROM documents ORDER BY download_count DESC LIMIT 10", fetch="all")],
        "top_users":     [dict(r) for r in execute_query("SELECT u.id,u.name,u.email,COUNT(d.id) as doc_count FROM users u LEFT JOIN documents d ON d.uploader_id=u.id GROUP BY u.id ORDER BY doc_count DESC LIMIT 10", fetch="all")],
        "daily_searches":[dict(r) for r in execute_query("SELECT DATE(created_at) as date, COUNT(*) as searches FROM search_logs GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30", fetch="all")],
    })

@app.route("/api/admin/activity")
@require_admin
def activity():
    rows = execute_query(
        "SELECT a.*, u.name as user_name, d.title as doc_title FROM activity_log a "
        "LEFT JOIN users u ON a.user_id=u.id LEFT JOIN documents d ON a.document_id=d.id "
        "ORDER BY a.created_at DESC LIMIT 100", fetch="all"
    )
    return jsonify({"activity":[dict(r) for r in rows]})

if __name__ == "__main__":
    init_db()
    log.info(f"GSU Repo starting on port {PORT} | AI={'ST' if HAS_ST else 'TFIDF'}")
    app.run(host="0.0.0.0", port=PORT, debug=False)
